import os
import io

try:
    from pypdf import PdfReader
except:
    PdfReader = None

try:
    import pandas as pd
except:
    pd = None

try:
    from docx import Document
except:
    Document = None

try:
    from pptx import Presentation
except:
    Presentation = None

try:
    from PIL import Image
    import pytesseract
except:
    Image = None
    pytesseract = None


def read_pdf(file_bytes):
    if not PdfReader: return ""
    try:
        reader = PdfReader(file_bytes)
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except:
        return ""


def read_text_bytes(data):
    return data.decode("utf-8", errors="ignore")


def read_excel(file_bytes):
    if not pd: return ""
    try:
        xls = pd.read_excel(file_bytes, sheet_name=None)
        rows = []
        for name, df in xls.items():
            rows.append(f"=== Sheet: {name} ===")
            rows += df.astype(str).apply(lambda r: " | ".join(r.values.tolist()), axis=1).tolist()
        return "\n".join(rows)
    except:
        return ""


def read_docx(file_bytes):
    if not Document: return ""
    try:
        doc = Document(file_bytes)
        return "\n".join(p.text for p in doc.paragraphs)
    except:
        return ""


def read_pptx(file_bytes):
    if not Presentation: return ""
    try:
        prs = Presentation(file_bytes)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    except:
        return ""


def read_image_ocr(file_bytes):
    if not Image or not pytesseract: return ""
    try:
        img = Image.open(file_bytes).convert("RGB")
        return pytesseract.image_to_string(img)
    except:
        return ""


def extract_text_from_upload(upload):
    name = upload.name
    data = upload.read()
    ext = os.path.splitext(name.lower())[1]

    if ext == ".pdf":
        return read_pdf(io.BytesIO(data)), name
    elif ext in [".txt", ".md", ".log", ".json", ".c", ".h", ".py"]:
        return read_text_bytes(data), name
    elif ext in [".xlsx", ".xls"]:
        return read_excel(io.BytesIO(data)), name
    elif ext == ".docx":
        return read_docx(io.BytesIO(data)), name
    elif ext == ".pptx":
        return read_pptx(io.BytesIO(data)), name
    elif ext in [".jpg", ".jpeg", ".png"]:
        return read_image_ocr(io.BytesIO(data)), name
    else:
        return read_text_bytes(data), name
