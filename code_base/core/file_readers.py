import io
import os
import re
import json
import requests
import pandas as pd
import traceback
from typing import List, Tuple
from pypdf import PdfReader
from core.config import MAX_CHARS_PER_CHUNK
from openpyxl import load_workbook

import email
from email import policy
from email.parser import BytesParser

import docx
from pptx import Presentation

try:
    from PIL import Image
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False


def split_text(text: str, max_chars: int = MAX_CHARS_PER_CHUNK) -> List[str]:
    paras = [p.strip() for p in text.split("\n\n") if p.strip()]
    chunks, buf = [], []
    size = 0
    for p in paras:
        if size + len(p) > max_chars and buf:
            chunks.append("\n\n".join(buf))
            buf, size = [p], len(p)
        else:
            buf.append(p)
            size += len(p)
    if buf:
        chunks.append("\n\n".join(buf))
    return chunks


# ---------------------------------------------------------
# File Readers
# ---------------------------------------------------------

def read_pdf(file: io.BytesIO) -> str:
    reader = PdfReader(file)
    texts = []
    for page in reader.pages:
        try:
            texts.append(page.extract_text() or "")
        except Exception:
            continue
    return "\n".join(texts)

def read_docx(file: io.BytesIO) -> str:
    doc = docx.Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def read_pptx(file: io.BytesIO) -> str:
    prs = Presentation(file)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def read_excel(file: io.BytesIO) -> str:
    err_pandas = err_openpyxl = err_html = err_csv = "Skipped or Unknown"
    file_bytes = file.read()
    try:
        dfs = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None,
                            dtype=str, keep_default_na=False)
        return _dfs_to_string(dfs)
    except Exception as e:
        err_pandas = str(e)
    try:
        wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        return _wb_to_string(wb)
    except Exception as e:
        err_openpyxl = str(e)
    try:
        dfs_list = pd.read_html(io.BytesIO(file_bytes))
        return _dfs_to_string({f"Table_{i}": df for i, df in enumerate(dfs_list)})
    except Exception as e:
        err_html = str(e)
    try:
        df = pd.read_csv(io.BytesIO(file_bytes), sep=None, engine='python')
        return _dfs_to_string({"CSV_Content": df})
    except Exception as e:
        err_csv = str(e)
    return (f"### EXCEL PROCESSING FAILED ###\n1. Pandas: {err_pandas}\n"
            f"2. OpenPyXL: {err_openpyxl}\n3. HTML: {err_html}\n4. CSV: {err_csv}")

def _dfs_to_string(dfs_dict):
    text_content = []
    for sheet_name, df in dfs_dict.items():
        row_count, col_count = len(df), len(df.columns)
        col_names = ", ".join(str(c) for c in df.columns)
        text_content.append(
            f"--- Sheet/Table: {sheet_name} ---\n"
            f"METADATA: This sheet has {row_count} data rows and {col_count} columns.\n"
            f"Columns: {col_names}\n--- Data Below ---"
        )
        text_content.append(df.fillna("").to_string(index=True, max_rows=None, max_cols=None))
    return "\n\n".join(text_content)

def _wb_to_string(wb):
    fallback_text = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        fallback_text.append(f"--- Sheet: {sheet_name} ---")
        rows = []
        for row in ws.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                rows.append(" | ".join(str(c) if c is not None else "" for c in row))
        fallback_text.append("\n".join(rows))
    return "\n\n".join(fallback_text)

def read_eml(file: io.BytesIO) -> str:
    try:
        msg = BytesParser(policy=policy.default).parse(file)
        parts = []
        for h in ['Subject', 'From', 'To', 'Date']:
            val = msg.get(h)
            if val:
                parts.append(f"{h}: {val}")
        parts.append("-" * 20)
        body = msg.get_body(preferencelist=('plain', 'html'))
        if body:
            try:
                parts.append(body.get_content())
            except Exception:
                parts.append("[Error decoding email body]")
        else:
            parts.append("[No text body found]")
        return "\n".join(parts)
    except Exception as e:
        return f"Error reading EML file: {str(e)}"

def read_image(file: io.BytesIO) -> str:
    if not OCR_AVAILABLE:
        return "[Error: PIL or Pytesseract not installed]"
    try:
        image = Image.open(file)
        text = pytesseract.image_to_string(image)
        return text if text.strip() else "[Image contained no readable text]"
    except Exception as e:
        return f"[Error: Tesseract OCR not found. Details: {str(e)}]"

def read_code_or_text(file: io.BytesIO) -> str:
    try:
        return file.read().decode("utf-8", errors="ignore")
    except Exception:
        return "[Binary or unreadable file content]"


# ---------------------------------------------------------
# Main Router
# ---------------------------------------------------------

def extract_text_from_upload(upload) -> Tuple[str, str, bytes]:
    name = upload.name
    data = upload.read()
    file_stream = io.BytesIO(data)
    ext = os.path.splitext(name.lower())[1]

    if ext in [".xlsx", ".xls"]:
        return read_excel(file_stream), name, data
    if ext == ".pdf":
        return read_pdf(file_stream), name, b""
    elif ext == ".docx":
        return read_docx(file_stream), name, b""
    elif ext == ".pptx":
        return read_pptx(file_stream), name, b""
    elif ext == ".eml":
        return read_eml(file_stream), name, b""
    elif ext in [".jpg", ".jpeg", ".png", ".bmp"]:
        return read_image(file_stream), name, b""
    elif ext in [".c", ".h", ".py", ".md", ".txt", ".json",
                 ".log", ".js", ".html", ".css"]:
        return read_code_or_text(file_stream), name, b""
    else:
        return read_code_or_text(file_stream), name, b""


# ---------------------------------------------------------
# Universal Excel Smart Filter
# ---------------------------------------------------------

# Groq TPM limit: 12,000 tokens/min on free tier
# 1 token â‰ˆ 4 chars. Reserve ~4000 tokens for system+question+answer.
# Remaining ~8000 tokens = ~32,000 chars for data context.
GROQ_SAFE_CONTEXT_CHARS = 30000


def _build_schema_summary(df: pd.DataFrame) -> str:
    """
    Build a compact schema: column name, inferred type, sample unique values.
    Sent to LLM so it knows what columns exist and what values look like.
    """
    lines = []
    for col in df.columns:
        sample_vals = df[col].replace("", pd.NA).dropna().head(20).tolist()
        numeric_count = sum(
            1 for v in sample_vals
            if re.match(r'^-?\d+(\.\d+)?$', str(v).strip())
        )
        if len(sample_vals) == 0:
            col_type = "empty"
        elif numeric_count == len(sample_vals):
            col_type = "numeric"
        elif numeric_count > len(sample_vals) / 2:
            col_type = "mostly_numeric"
        else:
            col_type = "text"

        unique_samples = list(dict.fromkeys(str(v) for v in sample_vals))[:5]
        lines.append(f'  - "{col}" ({col_type}): e.g. {unique_samples}')
    return "\n".join(lines)


def _ask_llm_for_filter_code(schema: str, query: str, api_key: str) -> str:
    """
    Ask the Groq LLM to generate a pandas boolean mask expression.
    The expression is later eval'd safely against the real DataFrame.
    """
    system_prompt = """You are a Python/pandas expert.
You will receive a DataFrame schema and a user query.
Write a single pandas boolean mask expression to filter the DataFrame.

STRICT RULES:
- DataFrame variable name: df
- All columns are dtype=str (loaded with dtype=str)
- For numeric comparisons: pd.to_numeric(df["ColName"], errors='coerce')
- For text search: df["ColName"].str.lower().str.contains("value", na=False)
- Combine with & (AND), | (OR). Always use parentheses around each condition.
- Output ONLY the expression. No imports, no assignments, no explanation, no markdown.
- If no filter needed (e.g. "show everything"), output: pd.Series([True] * len(df))

EXAMPLE:
Schema: "Age" (numeric), "City" (text)
Query: people older than 30 in London
Output: (pd.to_numeric(df["Age"], errors='coerce') > 30) & (df["City"].str.lower().str.contains("london", na=False))
"""

    user_prompt = f"""DataFrame schema:
{schema}

User query: "{query}"

Pandas boolean mask expression:"""

    try:
        response = requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json"
            },
            json={
                "model": "llama-3.3-70b-versatile",
                "max_tokens": 300,
                "temperature": 0,
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user",   "content": user_prompt}
                ]
            },
            timeout=30
        )
        data = response.json()
        raw = data["choices"][0]["message"]["content"].strip()
        # Strip markdown fences if present
        raw = re.sub(r'^```(?:python)?\s*', '', raw, flags=re.MULTILINE)
        raw = re.sub(r'\s*```$', '', raw, flags=re.MULTILINE)
        print(f"[DEBUG] LLM filter expression: {raw.strip()}")
        return raw.strip()
    except Exception as e:
        print(f"[DEBUG] LLM filter generation failed: {e}")
        return None


def _safe_eval_mask(expr: str, df: pd.DataFrame):
    """
    Safely
 evaluate a pandas mask expression.
    Blocks all Python builtins â€” only pd and df are accessible.
    """
    try:
        allowed_globals = {"pd": pd, "df": df, "__builtins__": {}}
        mask = eval(expr, allowed_globals)  # noqa
        if hasattr(mask, '__len__') and len(mask) == len(df):
            return mask.fillna(False).astype(bool)
        return None
    except Exception as e:
        print(f"[DEBUG] Mask eval error: {e}")
        return None


def _truncate_to_fit(df: pd.DataFrame, max_chars: int) -> Tuple[pd.DataFrame, bool]:
    """
    Truncate DataFrame rows so the string representation fits within max_chars.
    Returns (truncated_df, was_truncated).
    """
    if len(df) == 0:
        return df, False

    # Estimate chars per row from first 10 rows
    sample_str = df.head(10).to_string(index=False)
    chars_per_row = max(len(sample_str) / min(10, len(df)), 1)
    max_rows = max(1, int(max_chars / chars_per_row))

    if len(df) <= max_rows:
        return df, False

    return df.head(max_rows), True


def filter_excel_by_criteria(file_bytes: bytes, query: str,
                              api_key: str = None) -> str:
    """
    Universal Excel filter for ANY Excel file with ANY data.

    Pipeline:
    1. Load Excel â†’ build schema (column names + sample values)
    2. Send schema + query to Groq LLM â†’ get pandas filter expression
    3. Execute filter safely on DataFrame â†’ get matching rows
    4. Truncate result to fit Groq's token limit (30,000 chars)
    5. Return formatted text ready for LLM context

    Falls back to full data (truncated) if LLM filter step fails.
    """
    try:
        dfs = pd.read_excel(
            io.BytesIO(file_bytes),
            sheet_name=None,
            dtype=str,
            keep_default_na=False
        )
        results = []

        for sheet_name, df in dfs.items():
            df = df.fillna("").reset_index(drop=True)
            total_rows = len(df)
            col_names  = df.columns.tolist()

            print(f"[DEBUG] Sheet: {sheet_name} | Rows: {total_rows} | Cols: {len(col_names)}")

            # â”€â”€ Step 1: Build schema â”€â”€
            schema = _build_schema_summary(df)
            print(f"[DEBUG] Schema:\n{schema}")

            filtered_df   = None
            filter_method = "none"

            # â”€â”€ Step 2 & 3: LLM generates + executes filter â”€â”€
            if api_key:
                expr = _ask_llm_for_filter_code(schema, query, api_key)
                if expr:
                    mask = _safe_eval_mask(expr, df)
                    if mask is not None:
                        filtered_df   = df[mask].reset_index(drop=True)
                        filter_method = "llm_pandas"
                        print(f"[DEBUG] LLM filter matched: {len(filtered_df)} rows")
                    else:
                        print("[DEBUG] Mask eval failed â€” using full data")
                else:
                    print("[DEBUG] No LLM expression â€” using full data")
            else:
                print("[DEBUG] No API key provided â€” using full data")

            # â”€â”€ Fallback: full data â”€â”€
            if filtered_df is None:
                filtered_df   = df.copy()
                filter_method = "full_data"

            match_count = len(filtered_df)

            # â”€â”€ Step 4: Truncate to fit Groq token limit â”€â”€
            filtered_df, was_truncated = _truncate_to_fit(filtered_df, GROQ_SAFE_CONTEXT_CHARS)
            shown_count = len(filtered_df)

            print(f"[DEBUG] âœ… Final: showing {shown_count}/{match_count} matched rows "
                  f"(total={total_rows}, method={filter_method}, truncated={was_truncated})")

            # â”€â”€ Step 5: Format output â”€â”€
            header = (
                f"--- Sheet: {sheet_name} ---\n"
                f"Total rows in file: {total_rows} | "
                f"Rows matching query: {match_count} | "
                f"Rows shown (token limit): {shown_count}\n"
                f"Columns: {', '.join(col_names)}\n"
            )

            if was_truncated:
                header += (
                    f"âš ï¸� Result truncated to {shown_count} rows to fit model token limit. "
                    f"All {match_count} matched rows exist in the file.\n"
                )

            if match_count == 0:
                header = (
                    f"--- Sheet: {sheet_name} ---\n"
                    f"âš ï¸� No rows matched your query out of {total_rows} total rows.\n"
                    f"Columns: {', '.join(col_names)}\n\n"
                    f"Sample data (first 5 rows for reference):\n"
                    f"{df.head(5).to_string(index=False)}"
                )
                results.append(header)
            else:
                results.append(
                    header +
                    "--- Matching Data ---\n" +
                    filtered_df.to_string(index=False, max_rows=None, max_cols=None)
                )

        return "\n\n".join(results)

    except Exception as e:
        return f"Error in filter_excel_by_criteria: {str(e)}\n{traceback.format_exc()}"
