# app_gemini.py
"""
Knowledge Retention & Training Bot (Gemini API)
Streamlit + ChromaDB (local vectors) + Gemini API (chat + embeddings)
Supports: PDF, TXT/MD/LOG/code, DOCX, PPTX, XLSX, JPG/PNG (OCR)
"""

import os
import io
import uuid
import time
import traceback
from typing import List, Dict, Tuple

import streamlit as st
import chromadb
from chromadb.utils import embedding_functions
from openai import OpenAI

# ---------------------------
# GEMINI API key
# ---------------------------
GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY")
if not GEMINI_API_KEY:
    try:
        GEMINI_API_KEY = st.secrets["GEMINI_API_KEY"]
    except Exception:
        GEMINI_API_KEY = None

if GEMINI_API_KEY:
    client = OpenAI(api_key=GEMINI_API_KEY)
else:
    st.error("❗ GEMINI_API_KEY not found. Set it in environment or .streamlit/secrets.toml")
    st.stop()

# Optional heavy deps for file reading; handled gracefully if not present
try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    import pandas as pd
except Exception:
    pd = None

try:
    from docx import Document
except Exception:
    Document = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from PIL import Image
    import pytesseract
except Exception:
    Image = None
    pytesseract = None

# ---------------------------
# CONFIG
# ---------------------------
APP_TITLE = "Knowledge Retention & Training Bot (Gemini)"
PERSIST_DIR = ".chroma"
COLLECTION_NAME = "company_knowledge"
DEFAULT_TOP_K = 4
MAX_CHARS_PER_CHUNK = 1000
MAX_TOTAL_CONTEXT_CHARS = 1200
SUMMARIZE_THRESHOLD = 600

# Gemini chat/embedding models
DEFAULT_CHAT_MODEL = "gemini-1.5-t"        # adjust based on availability
DEFAULT_EMBED_MODEL = "gemini-embedding-1"

# ---------------------------
# Utilities: file readers
# ---------------------------
def read_pdf(file_bytes: io.BytesIO) -> str:
    if PdfReader is None:
        return ""
    try:
        reader = PdfReader(file_bytes)
        texts = []
        for page in reader.pages:
            try:
                texts.append(page.extract_text() or "")
            except Exception:
                continue
        return "\n".join(texts)
    except Exception:
        return ""

def read_text_bytes(data: bytes) -> str:
    try:
        return data.decode("utf-8", errors="ignore")
    except Exception:
        return ""

def read_excel(file_bytes: io.BytesIO) -> str:
    if pd is None:
        return ""
    try:
        xls = pd.read_excel(file_bytes, sheet_name=None)
        rows = []
        for name, df in xls.items():
            rows.append(f"=== Sheet: {name} ===")
            rows += df.astype(str).apply(lambda r: " | ".join(r.values.tolist()), axis=1).tolist()
        return "\n".join(rows)
    except Exception:
        return ""

def read_docx(file_bytes: io.BytesIO) -> str:
    if Document is None:
        return ""
    try:
        doc = Document(file_bytes)
        return "\n".join([p.text for p in doc.paragraphs])
    except Exception:
        return ""

def read_pptx(file_bytes: io.BytesIO) -> str:
    if Presentation is None:
        return ""
    try:
        prs = Presentation(file_bytes)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
                except Exception:
                    continue
        return "\n".join(texts)
    except Exception:
        return ""

def read_image_ocr(file_bytes: io.BytesIO) -> str:
    if Image is None or pytesseract is None:
        return ""
    try:
        img = Image.open(file_bytes).convert("RGB")
        return pytesseract.image_to_string(img)
    except Exception:
        return ""

def extract_text_from_upload(upload) -> Tuple[str, str]:
    name = upload.name
    data = upload.read()
    ext = os.path.splitext(name.lower())[1]
    if ext == ".pdf":
        return read_pdf(io.BytesIO(data)), name
    elif ext in [".txt", ".md", ".log", ".json", ".c", ".h", ".py"]:
        return read_text_bytes(data), name
    elif ext in [".xlsx", ".xls"]:
        return read_excel(io.BytesIO(data)), name
    elif ext in [".docx"]:
        return read_docx(io.BytesIO(data)), name
    elif ext in [".pptx"]:
        return read_pptx(io.BytesIO(data)), name
    elif ext in [".jpg", ".jpeg", ".png", ".tiff", ".bmp"]:
        return read_image_ocr(io.BytesIO(data)), name
    else:
        return read_text_bytes(data), name

# ---------------------------
# Embedding support via Gemini API
# ---------------------------
@st.cache_resource(show_spinner=False)
def get_embedder_gemini():
    def gemini_embed(texts: List[str]) -> List[List[float]]:
        out = []
        batch_size = 16
        for i in range(0, len(texts), batch_size):
            batch = texts[i : i + batch_size]
            resp = client.embeddings.create(
                model=DEFAULT_EMBED_MODEL,
                input=batch
            )
            for item in resp.data:
                out.append(item["embedding"])
        return out

    class _GeminiEmbed(embedding_functions.EmbeddingFunction):
        def __init__(self):
            super().__init__()
        def __call__(self, texts: List[str]) -> List[List[float]]:
            return gemini_embed(texts)

    return "gemini", _GeminiEmbed()

# ---------------------------
# Chroma client + collection
# ---------------------------
@st.cache_resource(show_spinner=False)
def get_chroma_client(persist_dir: str = PERSIST_DIR):
    os.makedirs(persist_dir, exist_ok=True)
    return chromadb.PersistentClient(path=persist_dir)

def ensure_collection(client: chromadb.ClientAPI, embed_fn) -> chromadb.Collection:
    try:
        col = client.get_collection(COLLECTION_NAME)
    except Exception:
        col = client.create_collection(COLLECTION_NAME, embedding_function=embed_fn)
        return col
    try:
        col = client.get_collection(COLLECTION_NAME, embedding_function=embed_fn)
    except Exception:
        pass
    return col

def add_document(chroma_col: chromadb.Collection, text: str, meta: Dict):
    if not text or not text.strip():
        return 0
    chunks = []
    paras = [p.strip() for p in text.split("\n\n") if p.strip()]
    buf = []
    size = 0
    for p in paras:
        if size + len(p) > MAX_CHARS_PER_CHUNK and buf:
            chunks.append("\n\n".join(buf))
            buf = [p]; size = len(p)
        else:
            buf.append(p); size += len(p)
    if buf:
        chunks.append("\n\n".join(buf))
    ids = [str(uuid.uuid4()) for _ in chunks]
    chroma_col.add(documents=chunks, metadatas=[meta] * len(chunks), ids=ids)
    return len(chunks)

def retrieve(chroma_col: chromadb.Collection, query: str, top_k: int = DEFAULT_TOP_K):
    res = chroma_col.query(query_texts=[query], n_results=top_k, include=["documents", "metadatas", "distances"])
    docs = res.get("documents", [[]])[0]
    metas = res.get("metadatas", [[]])[0]
    dists = res.get("distances", [[]])[0]
    ids = [str(i) for i in range(len(docs))]
    return list(zip(docs, metas, dists, ids))

# ---------------------------
# Gemini Chat wrapper
# ---------------------------
def call_gemini_chat(system_prompt: str, user_prompt: str, model: str = DEFAULT_CHAT_MODEL, timeout: int = 60) -> str:
    if not GEMINI_API_KEY:
        return "❗ Gemini API not configured. Set GEMINI_API_KEY."

    try:
        response = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            temperature=0.2,
            max_tokens=512,
        )
        return response.choices[0].message.content.strip()

    except Exception as e:
        tb = traceback.format_exc()
        return f"❗ Gemini API error:\n{e}\n\n{tb}"

# ---------------------------
# Summarization & prompt builder
# ---------------------------
def summarize_chunk(chunk: str) -> str:
    if len(chunk) <= SUMMARIZE_THRESHOLD:
        return chunk
    system = "You are a helpful assistant. Summarize the following text into concise bullet points."
    user = f"Text:\n{chunk}"
    return call_gemini_chat(system, user)

def build_prompt(question: str, hits: List[Tuple[str, Dict, float, str]]) -> Tuple[str, str]:
    sources_block = []
    total_chars = 0
    for i, (doc, meta, dist, _id) in enumerate(hits, start=1):
        title = meta.get("source", f"Doc {_id[:8]}") if isinstance(meta, dict) else f"Doc {_id[:8]}"
        chunk_text = summarize_chunk(doc) if len(doc) > SUMMARIZE_THRESHOLD else doc
        if total_chars + len(chunk_text) > MAX_TOTAL_CONTEXT_CHARS:
            remaining = MAX_TOTAL_CONTEXT_CHARS - total_chars
            if remaining <= 0:
                break
            chunk_text = chunk_text[:remaining]
        sources_block.append(f"[S{i}] {title}\n{chunk_text}")
        total_chars += len(chunk_text)
        if total_chars >= MAX_TOTAL_CONTEXT_CHARS:
            break
    context = "\n\n".join(sources_block) if sources_block else "No context available."
    system = (
        "You are an internal company knowledge assistant. Answer clearly and concisely. Cite sources like [S1]. "
        "If the answer is not in the provided context, say you don't know."
    )
    user = f"Question: {question}\n\nContext:\n{context}"
    return system, user

# ---------------------------
# Streamlit UI
# ---------------------------
st.set_page_config(page_title=APP_TITLE, page_icon="📚", layout="wide")
st.title(APP_TITLE)
st.caption("Upload documents, build a local knowledge base, then ask questions. Uses Gemini API for generation.")

with st.sidebar:
    st.subheader("Settings")
    top_k = st.slider("Top-K Context Chunks", 1, 5, 3)
    st.markdown("---")
    st.subheader("Add Documents")
    uploads = st.file_uploader("Upload files (pdf, txt, docx, pptx, xlsx, jpg, png)", accept_multiple_files=True)
    ingest_btn = st.button("Ingest to Knowledge Base")
    clear_btn = st.button("Clear Knowledge Base (danger)", type="secondary")

# Initialize embedder & chroma
with st.spinner("Initializing embedder & vector DB..."):
    try:
        embed_source, embed_fn = get_embedder_gemini()
        client_chroma = get_chroma_client(PERSIST_DIR)
        collection = ensure_collection(client_chroma, embed_fn)
    except Exception as e:
        st.error(f"Initialization failed: {e}")
        st.stop()

if clear_btn:
    with st.spinner("Clearing vector store..."):
        try:
            client_chroma.delete_collection(COLLECTION_NAME)
        except Exception:
            pass
        collection = ensure_collection(client_chroma, embed_fn)
        st.success("Knowledge base cleared.")

if ingest_btn and uploads:
    added_total = 0
    for up in uploads:
        text, fname = extract_text_from_upload(up)
        if not text or not text.strip():
            st.warning(f"No text extracted from {fname}; skipping.")
            continue
        meta = {"source": fname}
        added = add_document(collection, text, meta)
        added_total += added
    st.success(f"Ingested files into {added_total} chunks.")

# Chat state
if "messages" not in st.session_state:
    st.session_state.messages = []

for m in st.session_state.messages:
    with st.chat_message(m["role"]):
        st.markdown(m["content"])

question = st.chat_input("Ask about your uploaded knowledge...")
if question:
    st.session_state.messages.append({"role": "user", "content": question})
    with st.chat_message("user"):
        st.markdown(question)

    with st.chat_message("assistant"):
        with st.spinner("Retrieving context and drafting answer..."):
            hits = retrieve(collection, question, top_k=top_k)
            if not hits or all(h[0] == [] for h in hits):
                st.warning("No relevant context found. Please upload documents or broaden your question.")
                answer = "No relevant context found."
            else:
                system, user = build_prompt(question, hits)
                answer = call_gemini_chat(system, user)
                st.markdown(answer)

                if hits:
                    with st.expander("Sources used"):
                        for i, (doc, meta, dist, _id) in enumerate(hits, start=1):
                            title = meta.get("source", f"Doc {_id[:8]}") if isinstance(meta, dict) else f"Doc {_id[:8]}"
                            st.markdown(f"**[S{i}]** {title} — distance: {dist:.4f}")
                            st.code(doc[:MAX_CHARS_PER_CHUNK] + ("…" if len(doc) > MAX_CHARS_PER_CHUNK else ""))

    st.session_state.messages.append({"role": "assistant", "content": answer})

st.markdown(
    """
---
**Notes**
- Provide GEMINI_API_KEY via environment variable or .streamlit/secrets.toml.
- Tesseract OCR binary required for image OCR (install separately).
"""
)
